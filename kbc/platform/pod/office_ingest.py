"""office_ingest — pre-render binary office sources to markdown the agent can Read.

The served box receives raw/ as ORIGINAL bytes. The agent's Read tool (and the
model) handle pdf / images / plain text natively, but zip-based office formats
(.pptx / .xlsx / .docx) are opaque binary. At raw-install time this renders each
into a sibling `<name>.md` (e.g. deck.pptx -> deck.pptx.md) so the agent reads
clean markdown; the ORIGINAL file is left in place for provenance.

Deliberately lightweight per-format extraction (python-pptx / openpyxl /
python-docx) — no OCR / vision. Embedded pictures are written to a sibling
`<name>.assets/` and referenced from the markdown, so a diagram carrying a
deck's actual conclusion reaches the model instead of vanishing; whatever
cannot be kept is NAMED there rather than dropped. Imports are lazy so this
module stays importable without
the optional deps (e.g. a dep-less unit host); an actual conversion needs them
installed — they are baked into the box image.
"""
from __future__ import annotations

import hashlib
import os
import posixpath
import shutil
import tempfile
import zipfile
from collections.abc import Iterable, Iterator
from pathlib import Path

OFFICE_EXTS = (".pptx", ".xlsx", ".docx")
DEFAULT_MAX_DERIVED_BYTES = 512 * 1024 * 1024
DEFAULT_MAX_ARCHIVE_UNPACKED_BYTES = 512 * 1024 * 1024
DEFAULT_MAX_ARCHIVE_FILES = 10_000

# Where each format keeps its embedded media inside the archive. docProps/ is
# deliberately absent: the thumbnail there is the file's own preview image, not
# anything the document says.
OFFICE_MEDIA_PREFIX = {".pptx": "ppt/media/", ".docx": "word/media/", ".xlsx": "xl/media/"}
# What the compile model can actually look at. A deck's EMF/WMF vector diagram
# or an embedded video is NAMED in the markdown rather than extracted — the
# model still learns that a picture is there, which is the part that was
# missing when images vanished without trace.
EXTRACTABLE_MEDIA_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
# A chart is neither text nor a picture: it is DATA, stored as XML that nothing
# here renders. Left alone it is the one part of an Office file that vanishes
# with no trace whatever — and in an operations deck the chart is frequently
# the conclusion. Both formats keep the numbers recoverable: a deck caches
# them inside the chart part, and a workbook's chart points at cells this
# module already renders into the table above it.
OFFICE_CHART_PREFIX = {".pptx": "ppt/charts/", ".docx": "word/charts/", ".xlsx": "xl/charts/"}
_CHART_NS = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
_DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


class OfficeIngestLimitExceeded(ValueError):
    """An Office source would exceed a declared compile-box resource budget."""


def _positive_env(name: str, default: int) -> int:
    value = int(os.environ.get(name, str(default)))
    if value <= 0:
        raise ValueError(f"{name} must be positive")
    return value


class MediaSink:
    """Keeps an Office file's embedded pictures beside it, as real image
    sources the compile model can open.

    They land in `<name>.assets/`, which is the layout the rest of the platform
    already understands: the batch planner binds `<x>.assets/*` to `<x>.md`,
    and `<x>.md` is exactly the sibling this module writes, so a deck's
    pictures stay in the same batch as the deck and the coverage ledger
    auto-attaches them to it. Nothing downstream needs to change.

    A sink with no directory keeps nothing and reports every picture as a
    marker instead — the shape `convert_file` uses, where there is no tree to
    write into.
    """

    def __init__(self, assets_dir: Path | None, remaining: int) -> None:
        self.assets_dir = assets_dir
        self.remaining = remaining
        self.written = 0
        self._by_digest: dict[str, str] = {}
        self._emitted: set[str] = set()
        self.charts_placed = 0

    def note_chart(self) -> None:
        """Count a chart the parser placed, so the sweep can report only the
        ones it could not reach rather than repeating every chart."""
        self.charts_placed += 1

    @staticmethod
    def _digest(blob: bytes) -> str:
        return hashlib.sha1(blob).hexdigest()

    def note_emitted(self, blob: bytes) -> None:
        """Record that a picture already has a line in the document, so the
        sweep that follows the parser does not list it a second time.

        Keyed on CONTENT, not on the archive member name: the name is only
        reachable through library internals that differ between versions,
        while the bytes are the same bytes either path is holding."""
        self._emitted.add(self._digest(blob))

    def already_emitted(self, blob: bytes) -> bool:
        return self._digest(blob) in self._emitted

    def unkept_reason(self, ext: str) -> str:
        """Why a picture could not be kept, in the owner's terms. A marker that
        misstates the reason is worse than a terse one: saying a PNG cannot be
        parsed would send someone looking for a decoder problem that does not
        exist, when the truth is that this call site has nowhere to write."""
        ext = ("." + ext.lstrip(".")).lower() if ext else ""
        if ext not in EXTRACTABLE_MEDIA_EXTS:
            return f"{ext or '未知'} 格式，编译器不解析"
        if self.assets_dir is None:
            return "本次转换未落盘，请看原文件"
        return ""

    def keep(self, blob: bytes, ext: str) -> str | None:
        """Persist one picture and return the relative path to reference, or
        None when it cannot be kept. Identical bytes are stored once: a logo
        repeated on forty slides is one media part in the archive and must not
        become forty files."""
        ext = ("." + ext.lstrip(".")).lower()
        if self.assets_dir is None or ext not in EXTRACTABLE_MEDIA_EXTS or not blob:
            return None
        digest = self._digest(blob)
        if digest in self._by_digest:
            return self._by_digest[digest]
        if self.written + len(blob) > self.remaining:
            raise OfficeIngestLimitExceeded(
                f"embedded media exceeds the remaining {self.remaining}-byte Office budget"
            )
        name = f"{digest[:12]}{ext}"
        self.assets_dir.mkdir(parents=True, exist_ok=True)
        target = self.assets_dir / name
        if not target.exists():
            fd, temp_name = tempfile.mkstemp(prefix=f".{name}.", suffix=".tmp", dir=self.assets_dir)
            try:
                with os.fdopen(fd, "wb") as handle:
                    handle.write(blob)
                os.replace(temp_name, target)
            finally:
                if os.path.exists(temp_name):
                    os.unlink(temp_name)
        self.written += len(blob)
        rel = f"{self.assets_dir.name}/{name}"
        self._by_digest[digest] = rel
        return rel


def _picture_line(rel: str | None, caption: str, unkept: str = "") -> str:
    """One picture's markdown: a real reference when it was kept, an explicit
    note when it was not. Silence is the one thing this must never produce —
    the model has to know a picture is there even when it cannot see it."""
    if rel:
        return f"![{caption}]({rel})"
    detail = f"（{unkept}）" if unkept else ""
    return f"> [未提取的嵌入媒体] {caption}{detail}"


def _chart_lines(title: str, categories: list[str], series: list[tuple[str, list[str]]],
                 heading: str) -> Iterator[str]:
    """A chart as the table it always was. Categories down the left, one column
    per series — the shape a reader (and the model) can actually use."""
    yield f"### {heading}" + (f"：{title}" if title else "")
    if not categories or not series:
        return
    yield from _table_lines(
        [["分类"] + [name or f"系列{i + 1}" for i, (name, _) in enumerate(series)]]
        + [
            [category] + [values[row] if row < len(values) else "" for _, values in series]
            for row, category in enumerate(categories)
        ]
    )


def _parse_chart_xml(blob: bytes) -> tuple[str, list[str], list[tuple[str, list[str]]], list[str]]:
    """(title, categories, series, source refs) from a chart part, using the
    stdlib only. A deck caches its numbers here; a workbook usually does not,
    and then the refs say which cells hold them."""
    import xml.etree.ElementTree as ET

    try:
        root = ET.fromstring(blob)
    except ET.ParseError:
        return "", [], [], []

    def points(node) -> list[str]:
        if node is None:
            return []
        ordered: dict[int, str] = {}
        for pt in node.iter(_CHART_NS + "pt"):
            value = pt.find(_CHART_NS + "v")
            if value is None:
                continue
            try:
                ordered[int(pt.get("idx", "0"))] = (value.text or "").strip()
            except ValueError:
                continue
        return [ordered[k] for k in sorted(ordered)]

    title = ""
    title_node = root.find(f".//{_CHART_NS}title")
    if title_node is not None:
        title = "".join(t.text or "" for t in title_node.iter(_DRAWING_NS + "t")).strip()

    categories: list[str] = []
    series: list[tuple[str, list[str]]] = []
    for ser in root.iter(_CHART_NS + "ser"):
        name_node = ser.find(_CHART_NS + "tx")
        name = ""
        if name_node is not None:
            cached = points(name_node)
            name = cached[0] if cached else "".join(
                (f.text or "") for f in name_node.iter(_CHART_NS + "f")).strip()
        values = points(ser.find(_CHART_NS + "val"))
        if not categories:
            categories = points(ser.find(_CHART_NS + "cat"))
        series.append((name, values))
    refs = [(f.text or "").strip() for f in root.iter(_CHART_NS + "f") if (f.text or "").strip()]
    return title, categories, series, refs


def _sweep_archive_charts(path: Path, placed: int) -> Iterator[str]:
    """Report the charts the parser did not place.

    For a workbook or a document there is no position to recover, so every
    chart lands here. For a deck the shape walk already emitted the ones it
    could reach; a chart nested inside a grouped shape is not walked, and the
    count below is what keeps that from being a silent omission."""
    suffix = path.suffix.lower()
    prefix = OFFICE_CHART_PREFIX.get(suffix)
    if prefix is None:
        return
    try:
        with zipfile.ZipFile(path) as archive:
            parts = sorted(
                (
                    info for info in archive.infolist()
                    if not info.is_dir() and info.filename.startswith(prefix)
                    and info.filename.endswith(".xml")
                ),
                key=lambda i: i.filename,
            )
            if suffix == ".pptx":
                missing = len(parts) - placed
                if missing > 0:
                    yield (f"> [未定位的图表] 本文件另有 {missing} 个图表未能定位到具体幻灯片"
                           "（可能嵌在组合图形内）")
                return
            for index, info in enumerate(parts, 1):
                title, categories, series, refs = _parse_chart_xml(archive.read(info))
                yield from _chart_lines(title, categories, series, f"图表 {index}")
                if not (categories and series) and refs:
                    yield ("> 该图表未缓存数值，数据取自 "
                           + "、".join(dict.fromkeys(refs))
                           + "（对应数值见上方表格）")
    except (OSError, zipfile.BadZipFile):
        return


def _sweep_archive_media(path: Path, sink: MediaSink) -> Iterator[str]:
    """Emit whatever media the parser did not already place.

    Some pictures carry no position we can recover — a floating image in a
    .docx, a picture on an .xlsx sheet (openpyxl does not load images in the
    read-only mode this module depends on), a picture nested in a grouped
    shape. Reading them straight out of the archive needs no library at all, so
    the content survives even where the position does not."""
    prefix = OFFICE_MEDIA_PREFIX.get(path.suffix.lower())
    if prefix is None:
        return
    try:
        with zipfile.ZipFile(path) as archive:
            members = [
                info for info in archive.infolist()
                if not info.is_dir() and info.filename.startswith(prefix)
            ]
            heading_written = False
            for info in sorted(members, key=lambda i: i.filename):
                name = posixpath.basename(info.filename)
                ext = posixpath.splitext(name)[1].lower()
                extractable = ext in EXTRACTABLE_MEDIA_EXTS
                blob = archive.read(info) if extractable else b""
                if extractable and sink.already_emitted(blob):
                    continue  # the parser already placed this one, with its position
                rel = sink.keep(blob, ext) if extractable else None
                if not heading_written:
                    yield "## 附图（文档未标明位置）"
                    heading_written = True
                yield _picture_line(rel, name, "" if rel else sink.unkept_reason(ext))
    except (OSError, zipfile.BadZipFile):
        # The parser above already reported a format-specific error for a
        # corrupt archive; do not turn it into a second, less useful one.
        return


def _table_lines(rows: Iterable[list[str]]) -> Iterator[str]:
    """Stream a row matrix as a GitHub markdown table (first row = header)."""
    iterator = iter(rows)
    try:
        header = next(iterator)
    except StopIteration:
        return
    width = len(header)

    def cells(r: list[str]) -> str:
        r = (r + [""] * width)[:width]  # pad/truncate to header width → valid table
        return "| " + " | ".join(c.replace("\n", " ").replace("|", "\\|").strip() for c in r) + " |"

    yield cells(header)
    yield "| " + " | ".join("---" for _ in range(width)) + " |"
    for row in iterator:
        yield cells(row)


def _render(lines: Iterable[str]) -> str:
    return "\n".join(lines).strip() + "\n"


def _pptx_lines(path: Path, sink: MediaSink) -> Iterator[str]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        yield f"## Slide {i}"
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                # python-pptx reads the cached series publicly, so a deck's
                # chart arrives as the numbers it plots rather than as nothing.
                chart = shape.chart
                title = chart.chart_title.text_frame.text.strip() if chart.has_title else ""
                categories: list[str] = []
                series: list[tuple[str, list[str]]] = []
                for plot in chart.plots:
                    if not categories:
                        categories = [str(c) for c in plot.categories]
                    for one in plot.series:
                        series.append((
                            str(one.name or ""),
                            ["" if v is None else _number(v) for v in one.values],
                        ))
                yield from _chart_lines(title, categories, series, f"幻灯片 {i} 图表")
                sink.note_chart()
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # python-pptx exposes the bytes publicly here, so a deck — the
                # format most likely to carry its conclusion in a diagram —
                # keeps its pictures WITH the slide they belong to.
                image = shape.image
                yield _picture_line(
                    sink.keep(image.blob, image.ext), f"幻灯片 {i} 配图",
                    sink.unkept_reason(image.ext),
                )
                sink.note_emitted(image.blob)
            elif shape.has_table:
                yield from _table_lines([[c.text.strip() for c in row.cells] for row in shape.table.rows])
            elif shape.has_text_frame and shape.text_frame.text.strip():
                yield shape.text_frame.text.strip()


def pptx_to_md(path: Path) -> str:
    return _render(_lines_with_media(path, MediaSink(None, 0)))


def _xlsx_lines(path: Path, sink: MediaSink) -> Iterator[str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    try:
        for ws in wb.worksheets:
            yield f"## Sheet: {ws.title}"
            rows = (
                ["" if cell is None else str(cell) for cell in row]
                for row in ws.iter_rows(values_only=True)
            )
            yield from _table_lines(row for row in rows if any(cell.strip() for cell in row))
    finally:
        wb.close()


def xlsx_to_md(path: Path) -> str:
    return _render(_lines_with_media(path, MediaSink(None, 0)))


def _docx_lines(path: Path, sink: MediaSink) -> Iterator[str]:
    from docx import Document

    doc = Document(str(path))
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name or ""
        if style == "Title":
            yield "# " + text
        elif style.startswith("Heading"):
            tail = style.split()[-1]
            yield "#" * min(int(tail) if tail.isdigit() else 2, 6) + " " + text
        else:
            yield text
    for table in doc.tables:
        yield from _table_lines([[c.text.strip() for c in row.cells] for row in table.rows])


def docx_to_md(path: Path) -> str:
    return _render(_lines_with_media(path, MediaSink(None, 0)))


_CONVERTERS = {".pptx": pptx_to_md, ".xlsx": xlsx_to_md, ".docx": docx_to_md}
_LINE_CONVERTERS = {".pptx": _pptx_lines, ".xlsx": _xlsx_lines, ".docx": _docx_lines}


def _number(value) -> str:
    """Chart values arrive as floats. Render a whole number as a whole number —
    '989' is what the deck showed, not '989.0'."""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _lines_with_media(path: Path, sink: MediaSink) -> Iterator[str]:
    """The document's own lines, then the charts and media it left behind."""
    converter = _LINE_CONVERTERS[path.suffix.lower()]
    yield from converter(path, sink)
    yield from _sweep_archive_charts(path, sink.charts_placed)
    yield from _sweep_archive_media(path, sink)


def convert_file(path: Path) -> str | None:
    """Render one office file to markdown, or None if it is not an office format.

    Pictures are reported but not extracted: there is no tree to write them
    into. convert_tree is the path that keeps them."""
    fn = _CONVERTERS.get(path.suffix.lower())
    return fn(path) if fn else None


def _validate_archive_budget(path: Path) -> None:
    max_bytes = _positive_env("KBC_MAX_OFFICE_ARCHIVE_UNPACKED_BYTES", DEFAULT_MAX_ARCHIVE_UNPACKED_BYTES)
    max_files = _positive_env("KBC_MAX_OFFICE_ARCHIVE_FILES", DEFAULT_MAX_ARCHIVE_FILES)
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            unpacked = sum(info.file_size for info in infos)
    except (OSError, zipfile.BadZipFile):
        # Preserve the existing per-file corrupt-input handling; the real parser
        # will report the format-specific error to convert_tree.
        return
    if len(infos) > max_files:
        raise OfficeIngestLimitExceeded(
            f"{path.name} contains {len(infos)} archive entries; limit is {max_files}"
        )
    if unpacked > max_bytes:
        raise OfficeIngestLimitExceeded(
            f"{path.name} expands to {unpacked} bytes; Office archive limit is {max_bytes}"
        )


def _write_bounded_markdown(path: Path, dest: Path, remaining: int) -> int:
    _validate_archive_budget(path)
    # `deck.pptx` -> `deck.pptx.md` + `deck.pptx.assets/`: the batch planner
    # resolves `<x>.assets/*` to `<x>.md`, so the pictures ride in the same
    # batch as the markdown that references them.
    assets_dir = dest.with_name(dest.name[: -len(".md")] + ".assets")
    sink = MediaSink(assets_dir, remaining)
    fd, temp_name = tempfile.mkstemp(prefix=f".{dest.name}.", suffix=".tmp", dir=dest.parent)
    written = 0
    lines = 0
    try:
        with os.fdopen(fd, "wb") as output:
            for line in _lines_with_media(path, sink):
                payload = (line + "\n").encode("utf-8")
                if written + len(payload) > remaining:
                    raise OfficeIngestLimitExceeded(
                        f"{path.name} derived markdown exceeds the remaining {remaining}-byte Office budget"
                    )
                output.write(payload)
                written += len(payload)
                lines += 1
                if written + sink.written > remaining:
                    raise OfficeIngestLimitExceeded(
                        f"{path.name} derived markdown and media exceed the remaining "
                        f"{remaining}-byte Office budget"
                    )
            if lines == 0:
                # Preserve the previous converter contract for a valid but
                # empty Office file: it produced a one-newline readable
                # sidecar, rather than making the supported source look opaque.
                if remaining < 1:
                    raise OfficeIngestLimitExceeded(
                        f"{path.name} derived markdown exceeds the remaining {remaining}-byte Office budget"
                    )
                output.write(b"\n")
                written = 1
        os.replace(temp_name, dest)
        return written + sink.written
    except BaseException:
        # The markdown never landed, so its pictures must not linger either —
        # an orphaned `<x>.assets/` would be an unaccounted source referenced
        # by nothing, and the next run would find a half-converted document.
        shutil.rmtree(assets_dir, ignore_errors=True)
        raise
    finally:
        if os.path.exists(temp_name):
            os.unlink(temp_name)


def convert_tree(root: str) -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    """Walk `root`, render each .pptx/.xlsx/.docx to a sibling `<name>.md`.

    Returns (converted, errors): converted = [(src_rel, md_rel)], errors =
    [(src_rel, message)]. A corrupt parser input is recorded and skipped so one
    bad deck does not sink the KB install. Declared resource-budget violations
    fail the atomic snapshot commit instead of silently leaving a supported
    source unreadable. Idempotent: an existing `<name>.md` is left untouched (a
    re-install of the same bundle is a no-op)."""
    r = Path(root)
    converted: list[tuple[str, str]] = []
    errors: list[tuple[str, str]] = []
    max_derived = _positive_env("KBC_MAX_OFFICE_DERIVED_BYTES", DEFAULT_MAX_DERIVED_BYTES)
    derived = 0
    if not r.is_dir():
        return converted, errors
    for f in sorted(r.rglob("*")):
        if not f.is_file() or f.suffix.lower() not in OFFICE_EXTS:
            continue
        rel = f.relative_to(r).as_posix()
        dest = f.with_name(f.name + ".md")  # deck.pptx -> deck.pptx.md
        if dest.exists():
            continue
        try:
            written = _write_bounded_markdown(f, dest, max_derived - derived)
        except OfficeIngestLimitExceeded:
            # Resource-budget violations are not corrupt-file anomalies. The
            # source type is otherwise supported, so fail the atomic snapshot
            # commit rather than silently compiling an unreadable binary.
            raise
        except Exception as e:  # any parser failure is per-file non-fatal (fail-open boundary)
            errors.append((rel, repr(e)))
            continue
        if written == 0:
            continue
        derived += written
        converted.append((rel, dest.relative_to(r).as_posix()))
    return converted, errors
