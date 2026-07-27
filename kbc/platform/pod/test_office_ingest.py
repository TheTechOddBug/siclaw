"""Tests for office_ingest (pptx/xlsx/docx → sibling markdown pre-render).

Needs python-pptx / openpyxl / python-docx — the same deps the box image bakes
in. Run: python test_office_ingest.py
"""

import tempfile
from pathlib import Path
from unittest import mock

import office_ingest


def _make_samples(root: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "GPU 选型"
    s.placeholders[1].text = "H100 vs A100"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "拓扑"
    tb = s2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
    tb.cell(0, 0).text = "区域"; tb.cell(0, 1).text = "节点"
    tb.cell(1, 0).text = "华东"; tb.cell(1, 1).text = "52"
    prs.save(str(root / "deck.pptx"))

    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "配额"
    ws.append(["团队", "GPU"]); ws.append(["train", "32"])
    wb.save(str(root / "sub" / "quota.xlsx"))  # nested → exercises rglob + rel path

    from docx import Document
    d = Document(); d.add_heading("手册", level=1); d.add_paragraph("正文一句。")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "现象"; t.rows[0].cells[1].text = "动作"
    d.save(str(root / "manual.docx"))


def test_convert_tree():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "sub").mkdir()
        (root / "notes.md").write_text("already markdown", "utf-8")  # non-office → untouched
        _make_samples(root)
        converted, errors = office_ingest.convert_tree(str(root))
        assert not errors, errors
        assert dict(converted) == {
            "deck.pptx": "deck.pptx.md",
            "manual.docx": "manual.docx.md",
            "sub/quota.xlsx": "sub/quota.xlsx.md",
        }, converted
        pptx_md = (root / "deck.pptx.md").read_text("utf-8")
        assert "## Slide 1" in pptx_md and "GPU 选型" in pptx_md
        assert "| 区域 | 节点 |" in pptx_md and "| 华东 | 52 |" in pptx_md  # table preserved
        xlsx_md = (root / "sub" / "quota.xlsx.md").read_text("utf-8")
        assert "## Sheet: 配额" in xlsx_md and "| train | 32 |" in xlsx_md
        docx_md = (root / "manual.docx.md").read_text("utf-8")
        assert docx_md.startswith("# 手册") and "正文一句。" in docx_md and "| 现象 | 动作 |" in docx_md
        assert (root / "notes.md").read_text("utf-8") == "already markdown"  # untouched
        # idempotent: a re-run (same bundle re-installed) converts nothing new
        again, _ = office_ingest.convert_tree(str(root))
        assert again == [], again
    print("OK  convert_tree (pptx/xlsx/docx → sibling md; tables, nested, non-office untouched, idempotent)")


def test_fail_open_on_corrupt():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "broken.pptx").write_bytes(b"not a real pptx")  # corrupt
        from docx import Document
        doc = Document(); doc.add_paragraph("fine"); doc.save(str(root / "ok.docx"))
        converted, errors = office_ingest.convert_tree(str(root))
        assert [s for s, _ in converted] == ["ok.docx"], converted   # valid one still rendered
        assert [s for s, _ in errors] == ["broken.pptx"], errors     # corrupt recorded, NOT raised
        assert (root / "ok.docx.md").exists() and not (root / "broken.pptx.md").exists()
    print("OK  fail-open (corrupt file recorded in errors, valid files still converted, never raises)")


def test_derived_markdown_budget_is_atomic():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        from docx import Document
        doc = Document(); doc.add_paragraph("content that exceeds a tiny derived budget")
        doc.save(str(root / "large.docx"))

        with mock.patch.dict("os.environ", {"KBC_MAX_OFFICE_DERIVED_BYTES": "16"}):
            try:
                office_ingest.convert_tree(str(root))
                raise AssertionError("expected OfficeIngestLimitExceeded")
            except office_ingest.OfficeIngestLimitExceeded as error:
                assert "derived markdown" in str(error), error

        assert not (root / "large.docx.md").exists()
        assert not list(root.glob(".large.docx.md.*.tmp"))
    print("OK  derived markdown budget fails explicitly without a partial sidecar")


def test_archive_expansion_budget_is_atomic():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        from docx import Document
        doc = Document(); doc.add_paragraph("small document")
        doc.save(str(root / "archive.docx"))

        with mock.patch.dict("os.environ", {"KBC_MAX_OFFICE_ARCHIVE_UNPACKED_BYTES": "1"}):
            try:
                office_ingest.convert_tree(str(root))
                raise AssertionError("expected OfficeIngestLimitExceeded")
            except office_ingest.OfficeIngestLimitExceeded as error:
                assert "expands to" in str(error), error

        assert not (root / "archive.docx.md").exists()
        assert not list(root.glob(".archive.docx.md.*.tmp"))
    print("OK  archive expansion budget fails explicitly without a partial sidecar")


def test_empty_office_file_keeps_readable_sidecar_contract():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        from docx import Document
        Document().save(str(root / "empty.docx"))

        converted, errors = office_ingest.convert_tree(str(root))

        assert not errors, errors
        assert converted == [("empty.docx", "empty.docx.md")], converted
        assert (root / "empty.docx.md").read_bytes() == b"\n"
    print("OK  empty Office input retains its one-newline readable sidecar")


def _png(width=4, height=4) -> bytes:
    """A valid PNG — python-docx parses headers strictly, so a stub will not do."""
    import struct
    import zlib

    def chunk(kind: bytes, data: bytes) -> bytes:
        body = kind + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + bytes([200, 60, 60] * width) for _ in range(height))
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


def _deck_with_pictures(root: Path, *pictures: bytes) -> Path:
    import io

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "验收结论"
    for i, blob in enumerate(pictures):
        slide.shapes.add_picture(io.BytesIO(blob), Inches(1 + 2 * i), Inches(2), Inches(1), Inches(1))
    path = root / "deck.pptx"
    prs.save(str(path))
    return path


def test_deck_pictures_are_kept_beside_their_slide():
    """A deck's conclusion is often only in a diagram. Those pictures used to
    vanish without a trace: the sibling markdown never mentioned them, and the
    .pptx left in raw/ is opaque to the agent's Read tool."""
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        png = _png()
        # The same picture twice: one media part in the archive, and it must
        # become neither two files nor a duplicate 附图 entry.
        _deck_with_pictures(root, png, png)

        converted, errors = office_ingest.convert_tree(str(root))
        assert not errors, errors
        assert ("deck.pptx", "deck.pptx.md") in converted, converted

        body = (root / "deck.pptx.md").read_text(encoding="utf-8")
        assets = sorted((root / "deck.pptx.assets").iterdir())
        assert len(assets) == 1, f"identical pictures must be stored once: {assets}"
        assert assets[0].read_bytes() == png
        assert body.count(f"![幻灯片 1 配图](deck.pptx.assets/{assets[0].name})") == 2, body
        assert "附图" not in body, body
    print("OK  deck pictures kept, positioned by slide, stored once")


def test_images_land_where_the_platform_already_looks():
    """`<x>.assets/*` resolves to `<x>.md` in the batch planner and counts as an
    auto-attachable media asset in the ledger. The sibling this module writes IS
    `<x>.md`, so the layout must stay exactly this — otherwise the pictures
    become unaccounted sources planned into somebody else's batch."""
    import batching
    import selfcheck

    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        _deck_with_pictures(root, _png())
        office_ingest.convert_tree(str(root))

        inventory = batching.scan_sources(root)
        paths = {i["path"] for i in inventory}
        image = next(p for p in paths if p.startswith("deck.pptx.assets/"))
        assert batching._asset_anchor(image, paths) == "deck.pptx.md"
        assert selfcheck.is_media_asset(image)
        family = next(f for f in batching.source_families(inventory)
                      if any(i["path"] == "deck.pptx.md" for i in f))
        assert image in {i["path"] for i in family}, "picture must ride with its markdown"
    print("OK  extracted pictures bind to their document in planner and ledger")


def test_unreadable_media_is_named_rather_than_dropped():
    """A vector diagram or a clip cannot be handed to the model, but the model
    must still learn it is there — and the note must not misstate why."""
    import zipfile

    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        deck = _deck_with_pictures(root, _png())
        with zipfile.ZipFile(deck, "a") as archive:
            archive.writestr("ppt/media/diagram1.emf", b"\x01\x00\x00\x00 vector")

        office_ingest.convert_tree(str(root))
        body = (root / "deck.pptx.md").read_text(encoding="utf-8")

        assert "> [未提取的嵌入媒体] diagram1.emf（.emf 格式，编译器不解析）" in body, body
        # The PNG beside it WAS kept, so nothing may claim PNG is unparseable.
        assert "png 格式，编译器不解析" not in body, body
        assert "deck.pptx.assets/" in body, body
    print("OK  unreadable media is named, with an accurate reason")


def test_media_counts_against_the_office_budget():
    """Pictures are bytes on disk exactly as the markdown is. A budget counting
    only text would wave through a deck of two hundred photographs."""
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        _deck_with_pictures(root, _png(64, 64))

        with mock.patch.dict("os.environ", {"KBC_MAX_OFFICE_DERIVED_BYTES": "200"}):
            try:
                office_ingest.convert_tree(str(root))
            except office_ingest.OfficeIngestLimitExceeded:
                pass
            else:
                raise AssertionError("media bytes must count against the derived budget")
        assert not (root / "deck.pptx.md").exists(), "markdown must not survive the failure"
        assert not (root / "deck.pptx.assets").exists(), "pictures must not linger either"
    print("OK  media counts against the budget, and a failure leaves nothing behind")


def test_deck_chart_arrives_as_its_numbers():
    """A chart is the one part of an Office file that used to vanish with no
    trace at all — not even a marker — and in an ops deck it is often the
    conclusion. python-pptx reads the cached series publicly, so the deck's
    figures land in the markdown as the table they always were."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "算力对比"
        data = CategoryChartData()
        data.categories = ["H100", "H200", "B200"]
        data.add_series("TFLOPS", (989, 1979, 2250))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4), data)
        prs.save(str(root / "deck.pptx"))

        office_ingest.convert_tree(str(root))
        body = (root / "deck.pptx.md").read_text(encoding="utf-8")

        assert "### 幻灯片 1 图表" in body, body
        assert "| 分类 | TFLOPS |" in body, body
        # Whole numbers must read as whole numbers: the deck showed 989.
        assert "| H100 | 989 |" in body, body
        assert "989.0" not in body, body
        assert "| B200 | 2250 |" in body, body
        # The sweep must not then claim a chart went unplaced.
        assert "未定位的图表" not in body, body
    print("OK  deck chart arrives as its numbers, positioned by slide")


def test_workbook_chart_points_at_the_rows_it_plots():
    """A workbook usually caches nothing in the chart part — it does not need
    to, because the cells are right there. Naming the ranges is what makes the
    table above the chart legible as its data."""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference

    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        wb = Workbook()
        ws = wb.active
        ws.title = "验收"
        for row in (["机型", "TFLOPS"], ["H100", 989], ["H200", 1979]):
            ws.append(row)
        chart = BarChart()
        chart.title = "各机型算力"
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=3), titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=3))
        ws.add_chart(chart, "E5")
        wb.save(str(root / "report.xlsx"))

        office_ingest.convert_tree(str(root))
        body = (root / "report.xlsx.md").read_text(encoding="utf-8")

        assert "### 图表 1：各机型算力" in body, body
        assert "'验收'!$B$2:$B$3" in body, body
        assert "对应数值见上方表格" in body, body
        # The cells themselves are still rendered, which is what the note points at.
        assert "| H100 | 989 |" in body, body
    print("OK  workbook chart names its source ranges beside the rendered cells")


def test_cell_picture_in_a_workbook_is_kept():
    """openpyxl loads no images in the read-only mode this module depends on,
    so a picture sitting in a cell is reached through the archive instead. It
    keeps the content; only the cell it was anchored to is lost."""
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        picture = root / "cell.png"
        picture.write_bytes(_png(8, 8))
        wb = Workbook()
        ws = wb.active
        ws.append(["机型", "TFLOPS"])
        ws.append(["H100", 989])
        ws.add_image(XLImage(str(picture)), "D2")
        wb.save(str(root / "report.xlsx"))
        picture.unlink()  # only the workbook remains, as in a real raw/ tree

        office_ingest.convert_tree(str(root))
        body = (root / "report.xlsx.md").read_text(encoding="utf-8")

        assets = list((root / "report.xlsx.assets").iterdir())
        assert len(assets) == 1, assets
        assert assets[0].read_bytes() == _png(8, 8)
        assert f"report.xlsx.assets/{assets[0].name}" in body, body
        assert "| H100 | 989 |" in body, "the sheet's own data must still be there"
    print("OK  a picture inside a workbook cell is kept, with its rows intact")


def main():
    test_convert_tree()
    test_fail_open_on_corrupt()
    test_derived_markdown_budget_is_atomic()
    test_archive_expansion_budget_is_atomic()
    test_empty_office_file_keeps_readable_sidecar_contract()
    test_deck_pictures_are_kept_beside_their_slide()
    test_images_land_where_the_platform_already_looks()
    test_unreadable_media_is_named_rather_than_dropped()
    test_media_counts_against_the_office_budget()
    test_deck_chart_arrives_as_its_numbers()
    test_workbook_chart_points_at_the_rows_it_plots()
    test_cell_picture_in_a_workbook_is_kept()
    print("ALL OK  test_office_ingest")


if __name__ == "__main__":
    main()
